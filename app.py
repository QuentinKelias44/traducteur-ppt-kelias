import streamlit as st
from pptx import Presentation
from google import genai
import io
import json
import os

st.set_page_config(page_title="Traducteur de Fiches Produits PPT", page_icon="🌐")

st.title("🌐 Traducteur de Fiches Produits PowerPoint")
st.write("Conserve la charte graphique, le design et les zones de texte modifiables !")

api_key = st.text_input("Entre ta clé d'API Google AI Studio :", type="password")
langue_cible = st.selectbox("Choisir la langue de traduction :", ["Anglais", "Espagnol"])

uploaded_file = st.file_uploader("Dépose ton fichier PowerPoint (.pptx) ici", type=["pptx"])

if uploaded_file and api_key:
    if st.button("🚀 Lancer la traduction"):
        try:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            client = genai.Client(api_key=api_key.strip())
            prs = Presentation(uploaded_file)
            
            # Étape 1 : Extraction
            status_text.text("🔍 Extraction du texte de la présentation...")
            progress_bar.progress(20)
            
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    text_runs.append(run)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.text.strip():
                                            text_runs.append(run)

            if not text_runs:
                progress_bar.empty()
                status_text.empty()
                st.warning("⚠️ Aucun texte à traduire n'a été trouvé dans ce fichier.")
            else:
                # Étape 2 : Préparation et Envoi à l'IA
                status_text.text("🤖 Traduction en cours avec Gemini... Veuillez patienter quelques secondes.")
                progress_bar.progress(50)
                
                original_texts = [r.text for r in text_runs]
                
                prompt = f"""Tu es un traducteur expert en matériel industriel, signalisation et BTP.
Traduis la liste de textes suivante en {langue_cible}.

CONSIGNES STRICTES :
- Conserve le ton technique, concis et professionnel.
- Ne traduis PAS les normes (CE, NF, WL9, PL3, etc.), les dimensions (mm, kg, °C), ni les noms de marque (KELIAS, Px3 Plus, etc.).
- Renvoie UNIQUEMENT un tableau JSON contenant les traductions dans le même ordre exact sous la forme : ["texte1", "texte2", ...]

Textes à traduire :
{json.dumps(original_texts, ensure_ascii=False)}"""

                # Envoi API
                response = client.models.generate_content(
                    model='gemini-2.0-flash',
                    contents=prompt
                )
                
                # Étape 3 : Traitement de la réponse et réinjection
                status_text.text("✍️ Réinjection des textes traduits dans le document...")
                progress_bar.progress(80)
                
                clean_response = response.text.strip()
                if clean_response.startswith("```json"):
                    clean_response = clean_response[7:]
                elif clean_response.startswith("```"):
                    clean_response = clean_response[3:]
                if clean_response.endswith("```"):
                    clean_response = clean_response[:-3]
                clean_response = clean_response.strip()

                translated_texts = json.loads(clean_response)

                for run, trad in zip(text_runs, translated_texts):
                    run.text = trad

                # Étape 4 : Génération du nom de fichier dynamique
                nom_base, _ = os.path.splitext(uploaded_file.name)
                suffixe = ".ES" if langue_cible == "Espagnol" else ".EN"
                nom_fichier_final = f"{nom_base}{suffixe}.pptx"

                output = io.BytesIO()
                prs.save(output)
                output.seek(0)
                
                progress_bar.progress(100)
                status_text.text("✅ Traduction terminée avec succès !")
                
                st.success("🎉 Le document a été intégralement traduit !")
                
                st.download_button(
                    label=f"📥 Télécharger {nom_fichier_final}",
                    data=output,
                    file_name=nom_fichier_final,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

        except Exception as e:
            if 'progress_bar' in locals():
                progress_bar.empty()
            if 'status_text' in locals():
                status_text.empty()
                
            error_msg = str(e)
            
            if "429" in error_msg or "RESOURCE_EXHAUSTED" in error_msg:
                st.error("⏳ **Quota journalier dépassé.** La limite gratuite de Google a été atteinte pour aujourd'hui. Réessaie demain ou utilise une autre clé API.")
            elif "API_KEY_INVALID" in error_msg or "400" in error_msg:
                st.error("🔑 **Clé API invalide.** Vérifie que tu as bien copié-collé ta clé Google AI Studio sans espace en trop.")
            elif "503" in error_msg or "UNAVAILABLE" in error_msg:
                st.error("🌐 **Serveur temporairement indisponible.** Les serveurs de Google sont en surcharge. Patiente 1 minute et relance.")
            else:
                st.error(f"❌ **Une erreur est survenue :** {error_msg}")
